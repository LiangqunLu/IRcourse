#!/usr/bin/python3

from urllib.request import urlopen
from bs4 import BeautifulSoup
import re
import PyPDF2
from pptx import Presentation
import os, sys

###part I download text from the website
url = "http://www.cs.memphis.edu/~vrus/teaching/ir-websearch/"

html = urlopen(url)

soup = BeautifulSoup(html, 'html.parser')

#texts = soup.getText()

#find all links

links = []

for one in soup.find_all('a', href=True):
    
    http = one['href']
    
    if 'mailto' not in http or http != '':
        
        links.append(http)
    
 
links = [one for one in links if one != '' and 'mailto' not in one and '.php' not in one and 'jaffairs' not in one]   

#one function to clean text
def text_clean(text):
    
    texts = text.lower()
    texts = re.sub(r'\W+', ' ', texts)
    texts = re.sub(r'\d+', ' ', texts)
    texts = re.sub(r"_", ' ', texts)
    texts = re.sub(r"\s+", ' ', texts)
    texts = texts.strip()
    
    return(texts)

##one function to get text for each link
def Get_text(onelink, url = url):
    
    if 'http://' not in onelink:
        http = url + onelink
    else:
        http = onelink;
        
    ###read html
    if '.html' in onelink or onelink.startswith("http://"):
        
        html = urlopen(http)
        soup = BeautifulSoup(html, 'html.parser')
        text = soup.getText()    

    ##read pdf files    
    elif '.pdf' in onelink:
        
        os.system('wget -P tmp ' + http)
        pdfname = http.split('/')[-1]
        
        pdfFileObj = open('tmp/' + pdfname, 'rb')
        pdfReader = PyPDF2.PdfFileReader(pdfFileObj)
        pageObj = pdfReader.getPage(0)
        text = pageObj.extractText()
    
    ##read ppt python-pptx
    elif '.ppt' in onelink:
        
        a = os.system('wget -P tmp ' + http)
        
        if a == 0:
            pptname = http.split('/')[-1]
            
            prs = Presentation("tmp/" + pptname+"x")
        
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    
                for shape in slide.shapes:
            
                    if shape.has_text_frame: 
            
                        print(shape.text)
                        text.append(shape.text)
                        
            tmp = " ".join(text)
            text = tmp
        else:
            text = ""

    ##read ppt python-pptx
    elif '.txt' in onelink:
        
        a = os.system('wget -P tmp ' + http)
        if a == 0:        
            text = urlopen(http).read().decode('iso-8859-1')

        else:
            text = ""
    
    if text!="":
    
        text = text_clean(text)
    
    return(text)
        
##create a dic to store unique word and the frequency from main page and linked pages

links = [url] + links

alltext = []

for one in links:
    
    #print(one)
    out = Get_text(onelink = one)
    alltext.append(out)


alltext = [one for one in alltext if one != ""]

print("Total docs in this statistics is :", len(alltext))

##unique word number in all these docs
word_count = list(set(" ".join(alltext).split()))
word_count.sort()

print("Total unique word is :", len(word_count))

ff = open('output.txt', 'w')

header = [ 'freq_' + str(one + 1) for one in range(len(alltext))]

header = "\t".join(header)

ff.write(header + "\n")

word_freq = dict()

for one in word_count:
    
    freqs = [a.split().count(one) for a in alltext]

    word_freq[one] = freqs
    
    line = [one] + freqs; line = "\t".join([str(one) for one in line])
    
    #print(line)
    
    ff.write(line + "\n")
    
ff.close()    


