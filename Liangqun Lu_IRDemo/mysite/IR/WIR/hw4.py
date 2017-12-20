#!/usr/bin/python3

from urllib.request import urlopen
from bs4 import BeautifulSoup
import re
import PyPDF2
from pptx import Presentation
import os, sys
from bs4 import Comment
from stemming.porter2 import stem
import pandas as pd

os.chdir("/home/llu/Code/WIR/")

def stopwords():
    
    url = 'http://www.cs.memphis.edu/~vrus/teaching/ir-websearch/papers/english.stopwords.txt'
    stops = urlopen(url).read().decode('iso-8859-1').split("\n")
    stops = [one.strip() for one in stops if one != '']
    
    return(stops)

def text_clean_hw4(text):
    
    texts = text.lower()
    
    #remove digits
    texts = re.sub(r'\d+', ' ', texts)
    texts = re.sub(r'%', ' ', texts)
    texts = re.sub(r"\s+", ' ', texts)    
    ##remove stop words
    stops = stopwords()
    stops = [one.lower() for one in stops]  

    texts = " ".join([one for one in texts.split() if one not in stops])
    texts = re.sub(r'\W+', ' ', texts)

    texts = [stem(one) for one in texts.split()]
    texts = list(set(texts))
    
    texts = " ".join(texts)
    
    return(texts)

def getText(onelink):
    
    url = "https://www.yahoo.com" + onelink
    html = urlopen(url)
    soup = BeautifulSoup(html, 'html.parser')    
    
    for script in soup(["script", "style"]):
        script.extract()

    comments = soup.findAll(text=lambda text:isinstance(text, Comment))
    [comment.extract() for comment in comments]

    # get text
    text = u' '.join(soup.findAll(text=True))
    
    text = text_clean_hw4(text)
    
    return(text)    

def Get_text(onelink, url):
    
    if 'http://' not in onelink:
        http = url + onelink
    else:
        http = onelink;
        
    ###read html
    if '.html' in onelink or onelink.startswith("http://"):
        
        html = urlopen(http)
        soup = BeautifulSoup(html, 'html.parser')
        text = soup.getText()    

    ##read pdf files    
    elif '.pdf' in onelink:
        
        os.system('wget -P tmp ' + http)
        pdfname = http.split('/')[-1]
        
        pdfFileObj = open('tmp/' + pdfname, 'rb')
        pdfReader = PyPDF2.PdfFileReader(pdfFileObj)
        pageObj = pdfReader.getPage(0)
        text = pageObj.extractText()
    
    ##read ppt python-pptx
    elif '.ppt' in onelink:
        
        a = os.system('wget -P tmp ' + http)
        
        if a == 0:
            pptname = http.split('/')[-1]
            
            prs = Presentation("tmp/" + pptname+"x")
        
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    
                for shape in slide.shapes:
            
                    if shape.has_text_frame: 
            
                        print(shape.text)
                        text.append(shape.text)
                        
            tmp = " ".join(text)
            text = tmp
        else:
            text = ""

    ##read ppt python-pptx
    elif '.txt' in onelink:
        
        a = os.system('wget -P tmp ' + http)
        if a == 0:        
            text = urlopen(http).read().decode('iso-8859-1')

        else:
            text = ""
    
    if text!="":
    
        #text = text_clean(text)
        text = text_clean_hw4(text)
    return(text)

    
##obtain 10 news from Yahoo

url = "https://www.yahoo.com/news/"
html = urlopen(url)
soup = BeautifulSoup(html, 'html.parser')

links = [one['href'] for one in soup.find_all('a', href=True)]

links = [one for one in links if one.startswith('/news/')]

links = [one for one in links if '/m/' not in one][1:11]

for one in links:
    
    filename = one.split('/')[-1]
    
    out = open('news/' + filename, 'w')
    
    text = getText(one)
    
    out.write(text)
    
    out.close
    
###add txt, pdf files from hw 3
url = "http://www.cs.memphis.edu/~vrus/teaching/ir-websearch/"
html = urlopen(url)
soup = BeautifulSoup(html, 'html.parser')

links = [one['href'] for one in soup.find_all('a', href=True) if 'mailto' not in one['href'] or one['href'] != '' ]

links = [one for one in links if one != '' and 'mailto' not in one and '.php' not in one and 'jaffairs' not in one and '.ppt' not in one]   

links = [url] + links
links = list(set(links))

for one in links:
    
    if one.endswith("/"):
        one = one[:-1]
        
    filename = one.split('/')[-1]
    filename = re.sub('#', '', filename)
    
    out = open('news/' + filename, 'w')
    
    text = Get_text(onelink = one, url = url)
    
    out.write(text)
    
    out.close


